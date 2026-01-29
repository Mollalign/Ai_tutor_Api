"""
PPTX Presentation Parser

Extracts text from Microsoft PowerPoint presentations using python-pptx.

PPTX Structure:
--------------
Similar to DOCX, PPTX files are ZIP archives containing:
- ppt/slides/slide1.xml, slide2.xml, etc.
- ppt/slideMasters/: Master slide templates
- ppt/media/: Images and media

PowerPoint Content:
------------------
- Each slide can have multiple shapes (text boxes, tables, etc.)
- Shapes contain paragraphs with runs of text
- Notes pages are separate from slide content

What We Extract:
---------------
- Text from all shapes on each slide
- Slide titles (when identifiable)
- Table content
- Speaker notes (optional, valuable for learning)

Slides as Pages:
---------------
Unlike DOCX, presentations have natural page breaks (slides).
Each slide becomes a PageContent in our output.
"""

import io
import logging
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches
from pptx.exc import PackageNotFoundError

from app.ai.parsers.base import (
    DocumentParser,
    ParsedDocument,
    PageContent,
    ParserType,
)

logger = logging.getLogger(__name__)


class PPTXParser(DocumentParser):
    """
    Parser for Microsoft PowerPoint presentations (.pptx).
    
    Extracts text from slides, shapes, tables, and notes.
    Each slide is treated as a separate page.
    """
    
    def __init__(self, include_notes: bool = True):
        """
        Initialize PPTX parser.
        
        Args:
            include_notes: Whether to include speaker notes in output.
                          Notes often contain valuable explanations!
        """
        self.include_notes = include_notes
    
    @property
    def supported_types(self) -> List[ParserType]:
        return [ParserType.PPTX]
    
    def parse(
        self,
        content: bytes,
        filename: Optional[str] = None
    ) -> ParsedDocument:
        """
        Parse PPTX and extract text.
        
        Args:
            content: Raw PPTX file bytes
            filename: Original filename for logging
        
        Returns:
            ParsedDocument with extracted content
        """
        filename = filename or "unknown.pptx"
        logger.info(f"Parsing PPTX: {filename} ({len(content)} bytes)")
        
        try:
            # Create file-like object
            pptx_file = io.BytesIO(content)
            
            # Parse presentation
            prs = Presentation(pptx_file)
            
            # Extract metadata
            metadata = self._extract_metadata(prs, filename)
            
            # Extract content from each slide
            pages = []
            all_text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                # Extract slide content
                slide_text, slide_title = self._extract_slide_content(slide)
                
                # Extract speaker notes
                notes_text = ""
                if self.include_notes:
                    notes_text = self._extract_notes(slide)
                
                # Combine slide content
                full_slide_text = slide_text
                if notes_text:
                    full_slide_text += f"\n\n[Speaker Notes]\n{notes_text}"
                
                full_slide_text = self._clean_text(full_slide_text)
                
                # Create page content
                page_content = PageContent(
                    page_number=slide_num,
                    text=full_slide_text,
                    headings=[slide_title] if slide_title else [],
                    metadata={
                        "has_notes": bool(notes_text),
                        "shape_count": len(slide.shapes),
                    }
                )
                pages.append(page_content)
                
                if full_slide_text:
                    all_text_parts.append(
                        f"--- Slide {slide_num}"
                        f"{': ' + slide_title if slide_title else ''} ---\n"
                        f"{full_slide_text}"
                    )
            
            # Combine all text
            full_text = '\n\n'.join(all_text_parts)
            
            logger.info(
                f"PPTX parsed successfully: {filename}, "
                f"{len(pages)} slides, {len(full_text)} characters"
            )
            
            return ParsedDocument(
                text=full_text,
                pages=pages,
                metadata=metadata,
                success=True
            )
            
        except PackageNotFoundError as e:
            error_msg = f"Invalid PPTX file: {e}"
            logger.error(f"PPTX parse error for {filename}: {error_msg}")
            return ParsedDocument.from_error(error_msg)
            
        except Exception as e:
            error_msg = f"Error parsing PPTX: {e}"
            logger.exception(f"PPTX parse error for {filename}")
            return ParsedDocument.from_error(error_msg)
    
    def _extract_metadata(
        self,
        prs: Presentation,
        filename: str
    ) -> dict:
        """Extract metadata from PPTX."""
        metadata = {
            "filename": filename,
            "file_type": "pptx",
            "slide_count": len(prs.slides),
        }
        
        try:
            core_props = prs.core_properties
            
            if core_props.title:
                metadata["title"] = core_props.title
            if core_props.author:
                metadata["author"] = core_props.author
            if core_props.subject:
                metadata["subject"] = core_props.subject
                
        except Exception as e:
            logger.debug(f"Could not extract PPTX metadata: {e}")
        
        # Fallback title
        if "title" not in metadata:
            title = filename.rsplit('.', 1)[0]
            title = title.replace('_', ' ').replace('-', ' ')
            metadata["title"] = title
        
        return metadata
    
    def _extract_slide_content(self, slide) -> tuple[str, Optional[str]]:
        """
        Extract text from a single slide.
        
        Args:
            slide: python-pptx Slide object
        
        Returns:
            Tuple of (slide_text, slide_title)
        """
        text_parts = []
        slide_title = None
        
        for shape in slide.shapes:
            # Check for title
            if shape.is_placeholder:
                try:
                    if shape.placeholder_format.type == 1:  # Title placeholder
                        if hasattr(shape, 'text') and shape.text:
                            slide_title = shape.text.strip()
                except:
                    pass
            
            # Extract text from shape
            if hasattr(shape, 'text') and shape.text:
                text_parts.append(shape.text.strip())
            
            # Extract table content
            if shape.has_table:
                table_text = self._extract_table(shape.table)
                if table_text:
                    text_parts.append(table_text)
        
        slide_text = '\n'.join(text_parts)
        return slide_text, slide_title
    
    def _extract_notes(self, slide) -> str:
        """Extract speaker notes from slide."""
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text = notes_slide.notes_text_frame.text
                return notes_text.strip()
        except Exception as e:
            logger.debug(f"Could not extract notes: {e}")
        
        return ""
    
    def _extract_table(self, table) -> str:
        """Convert a PPTX table to text."""
        rows = []
        
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace('\n', ' ')
                cells.append(cell_text)
            
            if any(cells):
                rows.append(' | '.join(cells))
        
        return '\n'.join(rows)